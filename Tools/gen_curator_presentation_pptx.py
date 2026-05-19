# -*- coding: utf-8 -*-
"""
Презентация для куратора в колледже — прогресс после показа v0.1-demo.
Формат как в «Восстановление Курсора/gen_presentation_pptx.py»: python-pptx + слоты под скриншоты.

Зависимости: pip install python-pptx pillow

Запуск из корня репозитория MyProject:
  py -3 Tools\\gen_curator_presentation_pptx.py
"""
from __future__ import annotations

import os
import sys
from typing import Optional

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt
except ImportError:
    print("Установите: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

_BASE = os.path.dirname(os.path.abspath(__file__))
_ROOT = os.path.normpath(os.path.join(_BASE, ".."))
SCREEN_DIR = os.path.join(_BASE, "presentation_curator_screenshots")
OUT = os.path.join(_ROOT, "Docs", "CURATOR_COLLEGE_PRESENTATION.pptx")
OUT_ALT = os.path.join(_ROOT, "Docs", "CURATOR_COLLEGE_PRESENTATION_новая.pptx")

SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)


def find_screen(basename: str) -> Optional[str]:
    for ext in (".png", ".PNG", ".jpg", ".jpeg", ".JPG", ".JPEG"):
        p = os.path.join(SCREEN_DIR, basename + ext)
        if os.path.isfile(p):
            return p
    return None


def ensure_placeholder_image(target_png_path: str, line1: str, line2: str = "") -> str:
    if os.path.isfile(target_png_path):
        return target_png_path
    os.makedirs(os.path.dirname(target_png_path), exist_ok=True)
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        return ""

    w, h = 960, 540
    img = Image.new("RGB", (w, h), (228, 228, 232))
    dr = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("arial.ttf", 22)
        font_small = ImageFont.truetype("arial.ttf", 16)
    except OSError:
        font = font_small = ImageFont.load_default()
    dr.rectangle([0, 0, w, 40], fill=(200, 200, 210))
    dr.text((16, 10), "Место для скриншота", fill=(40, 40, 50), font=font)
    y = 200
    for t in (line1, line2):
        if t:
            dr.text((40, y), t[:90], fill=(70, 70, 80), font=font_small)
            y += 28
    dr.text((40, h - 36), os.path.basename(target_png_path), fill=(120, 120, 130), font=font_small)
    img.save(target_png_path, "PNG")
    return target_png_path


def pick_image(basename: str, hint1: str, hint2: str = "") -> Optional[str]:
    found = find_screen(basename)
    if found:
        return found
    placeholder = os.path.join(SCREEN_DIR, basename + ".png")
    return ensure_placeholder_image(placeholder, hint1, hint2) or None


def blank_slide(prs: Presentation):
    try:
        return prs.slides.add_slide(prs.slide_layouts[6])
    except Exception:
        return prs.slides.add_slide(prs.slide_layouts[5])


def add_title_slide(prs: Presentation, title: str, subtitle: str, right_image: Optional[str] = None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1 and slide.placeholders[1]:
        slide.placeholders[1].text = subtitle
    if right_image and os.path.isfile(right_image):
        slide.shapes.add_picture(right_image, Inches(5.2), Inches(1.5), width=Inches(4.5))


def add_slide_title_image_bullets(
    prs: Presentation,
    title: str,
    image_path: Optional[str],
    bullets: list[str],
    caption: str = "",
    font_pt: int = 14,
):
    slide = blank_slide(prs)
    tb_title = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.75))
    tb_title.text_frame.text = title
    p0 = tb_title.text_frame.paragraphs[0]
    p0.font.size = Pt(22)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(32, 32, 48)

    img_left = Inches(0.35)
    img_top = Inches(1.15)
    img_w = Inches(4.65)
    img_h = Inches(3.65)

    if image_path and os.path.isfile(image_path):
        slide.shapes.add_picture(image_path, img_left, img_top, width=img_w, height=img_h)
    else:
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, img_left, img_top, img_w, img_h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(230, 230, 235)
        shp.line.color.rgb = RGBColor(180, 180, 190)
        tf = shp.text_frame
        tf.text = "Нет изображения"
        tf.paragraphs[0].font.size = Pt(12)

    if caption:
        cap = slide.shapes.add_textbox(img_left, Inches(4.85), img_w, Inches(0.45))
        cap.text_frame.text = caption
        cap.text_frame.paragraphs[0].font.size = Pt(9)
        cap.text_frame.paragraphs[0].font.italic = True
        cap.text_frame.paragraphs[0].font.color.rgb = RGBColor(90, 90, 100)

    tb_body = slide.shapes.add_textbox(Inches(5.2), Inches(1.15), Inches(4.4), Inches(5.9))
    tf = tb_body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        if not line.strip():
            para = tf.add_paragraph()
            para.text = " "
            continue
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.font.size = Pt(font_pt)
        para.level = 0
        para.space_after = Pt(3)


def add_slide_two_images_one_row(
    prs: Presentation,
    title: str,
    left_path: Optional[str],
    right_path: Optional[str],
    bullets_bottom: list[str],
    captions: tuple[str, str] = ("Unreal Engine 5.3", "Visual Studio"),
):
    slide = blank_slide(prs)
    tb_title = slide.shapes.add_textbox(Inches(0.4), Inches(0.22), Inches(9.2), Inches(0.7))
    tb_title.text_frame.text = title
    tb_title.text_frame.paragraphs[0].font.size = Pt(22)
    tb_title.text_frame.paragraphs[0].font.bold = True

    y = Inches(1.0)
    w_img = Inches(4.55)
    h_img = Inches(2.85)
    gap = Inches(0.35)

    for i, (path, cap) in enumerate([(left_path, captions[0]), (right_path, captions[1])]):
        x = Inches(0.35) + i * (w_img + gap)
        if path and os.path.isfile(path):
            slide.shapes.add_picture(path, x, y, width=w_img, height=h_img)
        else:
            shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w_img, h_img)
            shp.fill.solid()
            shp.fill.fore_color.rgb = RGBColor(235, 235, 240)
        ctb = slide.shapes.add_textbox(x, y + h_img + Inches(0.05), w_img, Inches(0.35))
        ctb.text_frame.text = cap
        ctb.text_frame.paragraphs[0].font.size = Pt(10)
        ctb.text_frame.paragraphs[0].font.italic = True

    tb = slide.shapes.add_textbox(Inches(0.4), Inches(4.55), Inches(9.2), Inches(2.7))
    tf = tb.text_frame
    for i, line in enumerate(bullets_bottom):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.font.size = Pt(13)
        para.space_after = Pt(2)


def add_title_only_slide(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1 and slide.placeholders[1]:
        slide.placeholders[1].text = subtitle


def build():
    os.makedirs(SCREEN_DIR, exist_ok=True)

    prs = Presentation()
    prs.slide_width = int(SLIDE_W)
    prs.slide_height = int(SLIDE_H)

    add_title_slide(
        prs,
        "HonorKitchen / MyProject",
        "Отчёт куратору колледжа: прогресс после демо v0.1\n"
        "Unreal Engine 5.3 · C++ / Blueprints\n"
        "Сипатров Владислав Фёдорович — подставь группу при необходимости",
        pick_image("01_title", "Титульный кадр / меню игры", "01_title.png"),
    )

    add_slide_title_image_bullets(
        prs,
        "Что было на прошлом показе (v0.1-demo)",
        pick_image("02_since_demo", "Кадр с прошлого показа", "02_since_demo.png"),
        [
            "Рабочий цикл: меню → раунд → завершение.",
            "Сильная сторона по отзыву: процедурная генерация уровня (кухня).",
            "Базовые враги, портал, предметы; зафиксированы задачи на AI и визуал.",
        ],
        caption="По желанию: скрин с прошлой версии или текущего меню.",
    )

    add_slide_title_image_bullets(
        prs,
        "Что сделано с прошлого показа",
        pick_image("05_gameplay", "Геймплей / генерация", "05_gameplay.png"),
        [
            "Генерация кухни: устойчивее топология и спавн, dev-метрики (seed, nav, спавны).",
            "AI: рефакторинг TomatoSaurusAIController — стабильность состояний, Nav/Direct, watchdog.",
            "Инвентарь: стаки до 3, хотбар колесом мыши, тест-выдача для отладки.",
            "Документация: 14-дневный план, регресс-чеклист, журнал багов, инструкция установщика.",
        ],
        caption="Скрин процедурной кухни или HUD.",
        font_pt=13,
    )

    add_slide_title_image_bullets(
        prs,
        "Что осталось сделать (честный план)",
        pick_image("07_enemies", "Враги / сцена", "07_enemies.png"),
        [
            "Враги: довести преследование в режиме Chase (на дистанции возможно замирание; Melee стабильнее).",
            "Визуал: модели врагов и более «живая» карта (мебель/пропсы вместо голого прототипа).",
            "Дистрибуция: Shipping-сборка + установщик Inno Setup (см. Docs/INSTALLER_BUILD.md).",
        ],
        caption="Скрин врага или проблемного места — по желанию.",
        font_pt=14,
    )

    ue = pick_image("03_ue_editor", "Окно Unreal Editor", "03_ue_editor.png")
    vs = pick_image("04_vs", "Visual Studio", "04_vs.png")
    add_slide_two_images_one_row(
        prs,
        "Инструменты",
        ue,
        vs,
        [
            "Движок: Unreal Engine 5.3.",
            "Код: C++ (генератор кухни, AI, геймплей) + Blueprints.",
            "Сборка: Visual Studio; контроль версий: Git.",
        ],
        ("Unreal Engine", "Visual Studio"),
    )

    add_slide_title_image_bullets(
        prs,
        "Демо для куратора (3–5 минут)",
        pick_image("06_inventory", "Хотбар / предметы", "06_inventory.png"),
        [
            "1) Запуск из установщика или папки WindowsNoEditor.",
            "2) Новая игра — кратко цель: батарейки, портал.",
            "3) Показать генерацию и один контакт с врагом.",
            "4) Выход в меню без зависания.",
        ],
        caption="Хотбар или цель на HUD.",
    )

    add_slide_title_image_bullets(
        prs,
        "Установочный пакет",
        None,
        [
            "Упаковка: File → Package Project → Windows (64-bit), конфигурация Shipping.",
            "Установщик: Inno Setup 6, скрипт Tools/InnoSetup_MyProject_Demo.iss.",
            "На машине куратора: Windows 10/11 x64; при необходимости — VC++ Redistributable x64.",
            "Подробно: Docs/INSTALLER_BUILD.md и Docs/README_DEMO.txt.",
        ],
        caption="",
        font_pt=14,
    )

    add_title_only_slide(
        prs,
        "Спасибо за внимание",
        "Вопросы?\n\n"
        "Материалы: Docs/CURATOR_COLLEGE_PRESENTATION.md (текст), этот файл .pptx",
    )

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    try:
        prs.save(OUT)
        print("OK:", OUT)
    except OSError:
        prs.save(OUT_ALT)
        print("OK (основной файл занят):", OUT_ALT)


if __name__ == "__main__":
    build()
